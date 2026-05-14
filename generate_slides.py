"""Generate NLP_Healthcare_Mykhailovskyi.pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ─────────────────────────────────────────────────────────────────
BLUE       = RGBColor(0x1A, 0x56, 0xA8)
BLUE_LIGHT = RGBColor(0xE8, 0xF0, 0xFB)
BLUE_MID   = RGBColor(0x27, 0xA6, 0xE0)
GREEN      = RGBColor(0x1A, 0x7A, 0x40)
GREEN_LIGHT= RGBColor(0xE6, 0xF7, 0xEE)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)
ORANGE_LIGHT=RGBColor(0xFF, 0xF3, 0xE0)
RED        = RGBColor(0xC0, 0x39, 0x2B)
RED_LIGHT  = RGBColor(0xFD, 0xEC, 0xEA)
GRAY       = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT = RGBColor(0xF5, 0xF8, 0xFE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT_LINE= RGBColor(0xD0, 0xE2, 0xF8)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h,
                 size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
                 wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_para(tf, text, size=13, bold=False, color=BLACK, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def slide_header(slide, title, subtitle=None):
    """Blue top bar + title."""
    add_rect(slide, 0, 0, W, Inches(0.07), fill=BLUE)
    txb = slide.shapes.add_textbox(Inches(0.6), Inches(0.18), W - Inches(1.2), Inches(0.55))
    txb.word_wrap = False
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = BLUE
    # separator line
    add_rect(slide, Inches(0.6), Inches(0.73), W - Inches(1.2), Inches(0.025), fill=BLUE_LIGHT)
    if subtitle:
        add_text_box(slide, subtitle, Inches(0.6), Inches(0.75), W - Inches(1.2), Inches(0.3),
                     size=11, color=GRAY)

def slide_footer(slide, n, total=16):
    add_rect(slide, 0, H - Inches(0.32), W, Inches(0.32), fill=GRAY_LIGHT)
    add_rect(slide, 0, H - Inches(0.33), W, Inches(0.01), fill=ACCENT_LINE)
    add_text_box(slide, "CSCI E-222 · NLP for Healthcare · A. Mykhailovskyi",
                 Inches(0.4), H - Inches(0.3), Inches(9), Inches(0.25),
                 size=9, color=GRAY)
    add_text_box(slide, f"{n} / {total}",
                 W - Inches(1.2), H - Inches(0.3), Inches(1.0), Inches(0.25),
                 size=9, color=GRAY, align=PP_ALIGN.RIGHT)

def callout(slide, text, l, t, w, h, accent=BLUE, bg=BLUE_LIGHT, size=12):
    add_rect(slide, l, t, Inches(0.045), h, fill=accent)
    add_rect(slide, l + Inches(0.045), t, w - Inches(0.045), h, fill=bg)
    txb = slide.shapes.add_textbox(l + Inches(0.12), t + Inches(0.08),
                                   w - Inches(0.22), h - Inches(0.12))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = BLACK
    return txb

def metric_box(slide, value, label, l, t, w=Inches(1.8), h=Inches(0.85),
               val_color=BLUE, bg=GRAY_LIGHT, border=ACCENT_LINE):
    add_rect(slide, l, t, w, h, fill=bg, line=border, line_w=Pt(1))
    add_text_box(slide, value, l, t + Inches(0.1), w, Inches(0.42),
                 size=24, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, l, t + Inches(0.50), w, Inches(0.32),
                 size=10, color=GRAY, align=PP_ALIGN.CENTER)

def bullet_list(slide, items, l, t, w, h, size=12.5, indent=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        prefix = "• " if indent else ""
        run = p.add_run()
        run.text = prefix + item
        run.font.size = Pt(size)
        run.font.color.rgb = BLACK
    return txb

def section_label(slide, text, l, t, color=BLUE):
    txb = slide.shapes.add_textbox(l, t, Inches(5), Inches(0.28))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = color
    return txb

def add_table(slide, headers, rows, l, t, w, h, col_widths=None):
    ncols = len(headers)
    nrows = len(rows) + 1
    tbl = slide.shapes.add_table(nrows, ncols, l, t, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    # Header row
    for j, hdr in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = hdr
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = WHITE
    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            if (i + 1) % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_LIGHT
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(11)
            run.font.color.rgb = BLACK
    return tbl


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
# gradient-ish background tint
add_rect(sl, 0, 0, W, H, fill=RGBColor(0xF0, 0xF6, 0xFF))
add_rect(sl, 0, 0, W, Inches(0.10), fill=BLUE)

add_text_box(sl, "CSCI E-222 · Foundations of Large Language Models · Final Project",
             Inches(0.8), Inches(0.55), W - Inches(1.6), Inches(0.35),
             size=11, color=GRAY)

add_text_box(sl, "NLP for Healthcare:\nMedical Text Analysis",
             Inches(0.8), Inches(1.1), W - Inches(1.6), Inches(1.7),
             size=38, bold=True, color=BLUE)

add_text_box(sl,
             "Fine-tuning & Comparing Encoder vs. Decoder LLMs\non Medical NER and Question Answering",
             Inches(0.8), Inches(2.9), W - Inches(1.6), Inches(0.8),
             size=17, color=RGBColor(0x44, 0x44, 0x44))

add_text_box(sl, "A. Mykhailovskyi  ·  May 2026  ·  Harvard Extension School",
             Inches(0.8), Inches(3.85), W - Inches(1.6), Inches(0.35),
             size=12, color=GRAY)

# Pill badges
pills = [("Bio_ClinicalBERT", BLUE, BLUE_LIGHT),
         ("MedGemma-4B", BLUE, BLUE_LIGHT),
         ("MedMentions NER", GREEN, GREEN_LIGHT),
         ("MedQA USMLE", GREEN, GREEN_LIGHT)]
px = Inches(0.8)
for label, fc, bg in pills:
    w_pill = Inches(1.7)
    add_rect(sl, px, Inches(4.4), w_pill, Inches(0.32), fill=bg, line=fc, line_w=Pt(0.75))
    add_text_box(sl, label, px + Inches(0.1), Inches(4.43), w_pill - Inches(0.2), Inches(0.26),
                 size=11, bold=True, color=fc, align=PP_ALIGN.CENTER)
    px += w_pill + Inches(0.2)

add_text_box(sl, "Code: huggingface.co/alexd063  ·  Tasks: NER · Multi-choice QA",
             Inches(0.8), H - Inches(0.55), W - Inches(1.6), Inches(0.3),
             size=10, color=GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Problem Statement")
slide_footer(sl, 2)

callout(sl,
    "How do encoder-only (BioClinicalBERT) and decoder-based (MedGemma-4B) "
    "architectures compare when fine-tuned on two core medical NLP tasks — "
    "NER and QA — under a constrained compute budget?",
    Inches(0.5), Inches(0.95), W - Inches(1.0), Inches(0.85), BLUE, BLUE_LIGHT, 13)

# Left col
section_label(sl, "Why It Matters", Inches(0.5), Inches(1.95))
bullet_list(sl, [
    "Clinical NLP enables automated chart review, coding & decision support",
    "Domain-specific models outperform general LLMs but trade-offs are unclear",
    "LoRA/QLoRA makes billion-parameter models accessible on a single GPU",
    "Two very different architectures → clear ablation study",
], Inches(0.5), Inches(2.28), Inches(5.8), Inches(1.6))

section_label(sl, "Scope Constraints", Inches(0.5), Inches(3.95))
bullet_list(sl, [
    "Single GPU — Google Colab Pro A100",
    "PEFT only (LoRA/QLoRA) for the 4B model",
    "Public, reproducible datasets only",
], Inches(0.5), Inches(4.28), Inches(5.8), Inches(1.0))

# Right col
callout(sl,
    "Task 1 — Named Entity Recognition (NER)\n"
    "Label each token in biomedical text with a UMLS semantic type "
    "(e.g., T038 = Biologic Function, T103 = Chemical)",
    Inches(6.6), Inches(1.95), Inches(6.2), Inches(1.0), BLUE, BLUE_LIGHT, 12)

callout(sl,
    "Task 2 — Medical Question Answering\n"
    "Select the correct answer (A–D) from USMLE-style clinical questions "
    "requiring multi-step medical reasoning",
    Inches(6.6), Inches(3.1), Inches(6.2), Inches(1.0), GREEN, GREEN_LIGHT, 12)

callout(sl,
    "⚠  MedMentions uses a zero-shot label protocol — some test-set entity "
    "types are unseen during training, making evaluation harder across all models.",
    Inches(6.6), Inches(4.25), Inches(6.2), Inches(0.8), ORANGE, ORANGE_LIGHT, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DATASETS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Datasets")
slide_footer(sl, 3)

section_label(sl, "Task 1 · NER — MedMentions MTI881", Inches(0.5), Inches(0.95))
add_table(sl,
    ["Split", "Samples", "Approx. Tokens"],
    [["Train", "23,399", "~700 K"],
     ["Validation", "2,924", "~88 K"],
     ["Test", "2,926", "~88 K"]],
    Inches(0.5), Inches(1.28), Inches(4.5), Inches(1.1),
    col_widths=[Inches(1.5), Inches(1.5), Inches(1.5)])

bullet_list(sl, [
    "4,392 PubMed abstracts, 350 K+ UMLS-annotated mentions",
    "105 entity classes (UMLS semantic types), BIO format",
    'Severe class imbalance: "O" tag = 76% of all tokens',
    "Source: Ben10x/MedMentions-MTI881-NER",
], Inches(0.5), Inches(2.55), Inches(5.8), Inches(1.5))

callout(sl,
    "⚠  Dataset has no stated license — legal status is ambiguous. "
    "Usage proceeds under research exemption but should carry a disclaimer.",
    Inches(0.5), Inches(4.15), Inches(6.0), Inches(0.68), ORANGE, ORANGE_LIGHT, 11)

section_label(sl, "Task 2 · QA — MedQA USMLE-4-options", Inches(6.8), Inches(0.95))
add_table(sl,
    ["Split", "Samples"],
    [["Train", "10,200"], ["Test", "1,270"]],
    Inches(6.8), Inches(1.28), Inches(3.2), Inches(0.82),
    col_widths=[Inches(1.6), Inches(1.6)])

bullet_list(sl, [
    "USMLE Step 1/2/3 clinical reasoning questions",
    "4-option multiple choice (A–D)",
    "Random baseline: 25%  ·  Human expert: ~90%",
    "License: CC-BY-4.0 (permissive)",
    "Source: GBaker/MedQA-USMLE-4-options",
    "Standard benchmark: MedGemma 4B scores 64.4% (pre-trained)",
], Inches(6.8), Inches(2.25), Inches(6.0), Inches(1.7))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — MODEL OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Models & Methods Overview")
slide_footer(sl, 4)

# Left: BERT
section_label(sl, "Bio_ClinicalBERT  (Encoder)", Inches(0.5), Inches(0.95))
bullet_list(sl, [
    "110 M parameters, BERT-base architecture",
    "Pre-trained on MIMIC-III clinical notes (~880 M words)",
    "Fine-tuned with full parameter update",
    "NER → Token Classification head",
    "QA  → Multiple Choice head (4 seqs → 1 logit each)",
    "Bidirectional attention — strong at classification",
    "Max sequence length: 128 tokens",
], Inches(0.5), Inches(1.28), Inches(5.8), Inches(2.4))

callout(sl, "Architecture: Input → [ClinicalBERT] → Label / Choice",
        Inches(0.5), Inches(3.85), Inches(5.8), Inches(0.48), BLUE, BLUE_LIGHT, 11)

callout(sl, "MIT License — permissive for research use.",
        Inches(0.5), Inches(4.42), Inches(5.8), Inches(0.38), GREEN, GREEN_LIGHT, 11)

# Right: MedGemma
section_label(sl, "MedGemma-4B-IT  (Decoder + QLoRA)", Inches(6.8), Inches(0.95))
bullet_list(sl, [
    "4.3 B parameters, Gemma-3 architecture, 128 K context",
    "Pre-trained on medical images + text (Google Health AI)",
    "PEFT: QLoRA (4-bit NF4 + LoRA rank 16, alpha 16)",
    "Only 0.76% of parameters trainable (~33 M)",
    "NER → generative: \"token: tag, token: tag …\"",
    "QA  → generative: outputs a single letter A / B / C / D",
    "2× training speed via Unsloth optimization",
], Inches(6.8), Inches(1.28), Inches(6.0), Inches(2.4))

# Mini metrics row
metric_box(sl, "33 M", "Trainable params", Inches(6.8), Inches(3.82), Inches(1.9), Inches(0.78), BLUE)
metric_box(sl, "4-bit", "NF4 quantization", Inches(6.8) + Inches(2.0), Inches(3.82), Inches(1.9), Inches(0.78), BLUE)
metric_box(sl, "2×", "Faster via Unsloth", Inches(6.8) + Inches(4.0), Inches(3.82), Inches(1.9), Inches(0.78), GREEN)

callout(sl, "Google Health AI Developer Foundations terms required — restricts commercial use.",
        Inches(6.8), Inches(4.72), Inches(6.0), Inches(0.38), ORANGE, ORANGE_LIGHT, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — NER PIPELINE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Task 1 — NER: Implementation Pipeline")
slide_footer(sl, 5)

# Step strip
steps = ["Load\nMedMentions", "Tokenize +\nalign labels",
         "Compute\nclass weights", "WeightedTrainer\nfine-tune",
         "seqeval F1\nevaluation", "Push to\nHF Hub"]
step_w = (W - Inches(1.0)) / len(steps)
for i, s in enumerate(steps):
    x = Inches(0.5) + i * step_w
    add_rect(sl, x + step_w * 0.35, Inches(1.05), step_w * 0.3, Inches(0.3),
             fill=BLUE, line=BLUE)
    add_text_box(sl, str(i + 1), x + step_w * 0.35, Inches(1.05), step_w * 0.3, Inches(0.3),
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(sl, s, x, Inches(1.42), step_w, Inches(0.42),
                 size=10, color=BLACK, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_rect(sl, x + step_w * 0.65, Inches(1.17), step_w * 0.35, Inches(0.05),
                 fill=ACCENT_LINE)

# Left col
section_label(sl, "Key Engineering Decisions", Inches(0.5), Inches(2.05))
bullet_list(sl, [
    "Subword tokens → -100 label (ignored in cross-entropy loss)",
    "Exponential class weights: exp(−nC/N)×100, O-class fixed at 0.5",
    "Custom WeightedTrainer injects per-class loss weights",
    "Cosine LR schedule with 10% warmup, LR = 6e-5",
    "Early stopping (patience = 2 on macro F1)",
    "Checkpoints: alexd063/bio-clinicalbert-finetuned-medmentions",
], Inches(0.5), Inches(2.38), Inches(5.8), Inches(2.3))

# Right col
section_label(sl, "Three Experiments", Inches(6.8), Inches(2.05))
add_table(sl,
    ["Notebook", "What changed", "Status"],
    [["01-baseline",  "Standard fine-tuning",           "✓ Complete"],
     ["02-weighted",  "+ Class weights, cosine LR",      "✓ Complete"],
     ["wip-crf",      "+ CRF layer (seq. modelling)",    "✗ Init error"]],
    Inches(6.8), Inches(2.38), Inches(6.0), Inches(1.18),
    col_widths=[Inches(1.8), Inches(2.5), Inches(1.7)])

callout(sl,
    "CRF notebook fails: _tied_weights_keys = [] should be {} "
    "causing AttributeError in HuggingFace post_init(). Fix is a one-line change.",
    Inches(6.8), Inches(3.65), Inches(6.0), Inches(0.65), ORANGE, ORANGE_LIGHT, 11)

callout(sl,
    "✓  Class weighting was the single most impactful engineering decision — "
    "without it the model collapses to predicting 'O' for all tokens.",
    Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.58), GREEN, GREEN_LIGHT, 12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MEDGEMMA NER
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Task 1 — NER: MedGemma-4B with QLoRA")
slide_footer(sl, 6)

section_label(sl, "Generative NER Approach", Inches(0.5), Inches(0.95))
callout(sl,
    "Unlike token classification, MedGemma generates a free-text sequence:\n"
    "\"aspirin: B-T103, prescribed: O, daily: O, …\"\n"
    "Each token is explicitly labelled in the output text.",
    Inches(0.5), Inches(1.25), Inches(6.0), Inches(0.92), BLUE, BLUE_LIGHT, 12)

section_label(sl, "Training Configuration", Inches(0.5), Inches(2.3))
bullet_list(sl, [
    "QLoRA: 4-bit NF4, rank 16, alpha 16, no dropout",
    "Unsloth for 2× training speed",
    "SFTTrainer, max 300 steps (≈ 20% of 1 epoch)",
    "LR = 2e-4, batch = 4, grad accum = 4 (eff. batch 16)",
    "Evaluation: qualitative inspection of 3 validation samples only",
    "Checkpoint: alexd063/bio-gemma4bit-finetuned-medmentions",
], Inches(0.5), Inches(2.63), Inches(5.8), Inches(2.0))

callout(sl,
    "⚠  Only 300 of ~1,462 steps trained (severe underfitting).\n"
    "No quantitative F1 — evaluation is visual only.\n"
    "Parsing 'token: tag' format from free text is fragile.",
    Inches(0.5), Inches(4.72), Inches(6.0), Inches(0.68), ORANGE, ORANGE_LIGHT, 11)

section_label(sl, "Prompt Format", Inches(6.8), Inches(0.95))
# Prompt box
add_rect(sl, Inches(6.8), Inches(1.25), Inches(6.0), Inches(2.5),
         fill=RGBColor(0xF5, 0xF8, 0xFF), line=ACCENT_LINE, line_w=Pt(1))
prompt_text = (
    "### Instruction:\n"
    "Extract medical entities and label\n"
    "each token with its NER tag.\n\n"
    "### Input:\n"
    "Patient presents with acute appendicitis\n"
    "and elevated CRP levels.\n\n"
    "### Response:\n"
    "Patient: O, presents: O, with: O,\n"
    "acute: B-T047, appendicitis: I-T047,\n"
    "and: O, elevated: O, CRP: B-T059, …"
)
add_text_box(sl, prompt_text, Inches(7.0), Inches(1.35), Inches(5.6), Inches(2.3),
             size=10.5, color=BLACK)

section_label(sl, "Key Limitation", Inches(6.8), Inches(3.88))
callout(sl,
    "Generative NER is promising for arbitrary schema but requires structured "
    "output enforcement (constrained decoding or JSON mode) to be reliable at scale.",
    Inches(6.8), Inches(4.18), Inches(6.0), Inches(0.75), RED, RED_LIGHT, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — QA PIPELINE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Task 2 — QA: Implementation Pipeline")
slide_footer(sl, 7)

# Left col
section_label(sl, "BioClinicalBERT Multiple-Choice", Inches(0.5), Inches(0.95))
callout(sl,
    "(Q + Option A) → BERT → logit  ×4 → softmax → argmax",
    Inches(0.5), Inches(1.25), Inches(5.8), Inches(0.38), BLUE, BLUE_LIGHT, 11)

bullet_list(sl, [
    "Custom DataCollatorForMultipleChoice",
    "3 metrics: Accuracy, Macro F1, ECE (calibration)",
    "LR = 4e-5, cosine schedule, 10 epochs max",
    "Truncation side = left (preserves answer text)",
    "Early stopping patience = 2",
    "90/10 train/val split with seed=42",
    "Checkpoint: alexd063/bio-clinicalbert-finetuned-medqa",
], Inches(0.5), Inches(1.72), Inches(5.8), Inches(2.4))

callout(sl,
    "ECE (Expected Calibration Error) is an excellent addition — measures "
    "whether model confidence matches actual accuracy, beyond just raw metrics.",
    Inches(0.5), Inches(4.22), Inches(5.8), Inches(0.65), GREEN, GREEN_LIGHT, 11)

# Right col
section_label(sl, "MedGemma Generative QA", Inches(6.8), Inches(0.95))
callout(sl,
    "Model sees full question + all 4 options and outputs a single letter.\n"
    "Instruction format: explicit EOS token appended to answer.",
    Inches(6.8), Inches(1.25), Inches(6.0), Inches(0.6), BLUE, BLUE_LIGHT, 11)

bullet_list(sl, [
    "Custom MedicalQAEvaluationCallback — accuracy mid-training",
    "Temperature = 0.1 at inference (near-deterministic)",
    "2 epochs, LR = 2e-4, LoRA rank 16 on all attention matrices",
    "Packing enabled for throughput",
    "10 warmup steps (0.87% of total — too low)",
    "Checkpoint: alexd063/bio-gemma4bit-finetuned-medqa",
], Inches(6.8), Inches(1.95), Inches(6.0), Inches(2.1))

callout(sl,
    "⚠  Only 2 epochs trained. Warmup = 10 steps is far too low for stable "
    "fine-tuning. No error analysis on incorrect predictions.",
    Inches(6.8), Inches(4.15), Inches(6.0), Inches(0.65), ORANGE, ORANGE_LIGHT, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — NER RESULTS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Results — Task 1: Named Entity Recognition")
slide_footer(sl, 8)

metric_box(sl, "0.65", "BioClinicalBERT\nMacro F1 (full test)", Inches(0.5), Inches(0.95), Inches(2.9), Inches(0.9), GREEN)
metric_box(sl, "0.27", "MedGemma\nF1 (10% subset)", Inches(3.55), Inches(0.95), Inches(2.9), Inches(0.9), ORANGE)
metric_box(sl, "105", "Entity classes\n(UMLS types)", Inches(6.6), Inches(0.95), Inches(2.9), Inches(0.9), BLUE)
metric_box(sl, "76%", "\"O\" tokens\n(class imbalance)", Inches(9.65), Inches(0.95), Inches(2.9), Inches(0.9), BLUE)

section_label(sl, "BioClinicalBERT — Entity Performance", Inches(0.5), Inches(2.05))
add_table(sl,
    ["Top performers", "F1", "Bottom performers", "F1"],
    [["T204 (Receptor)",    "High", "T059 (Lab Procedure)", "~0"],
     ["T103 (Chemical)",    "High", "T024 (Cell Component)", "~0"],
     ["T017 (Anatomy)",     "High", "T081 (Quantitative Concept)", "~0"]],
    Inches(0.5), Inches(2.38), Inches(5.8), Inches(1.2),
    col_widths=[Inches(2.0), Inches(0.7), Inches(2.3), Inches(0.8)])

callout(sl,
    "Class weighting helps frequent entities but very rare classes "
    "(< 5 training examples) remain near zero — a fundamental data sparsity problem.",
    Inches(0.5), Inches(3.72), Inches(5.8), Inches(0.65), BLUE, BLUE_LIGHT, 11)

callout(sl,
    "✓  Key finding: Encoder architecture is clearly better at token-level "
    "classification. Generative NER is promising but needs output enforcement.",
    Inches(0.5), Inches(4.47), Inches(5.8), Inches(0.65), GREEN, GREEN_LIGHT, 11)

section_label(sl, "Evaluation Caveat", Inches(6.8), Inches(2.05))
callout(sl,
    "Asymmetric evaluation: BioClinicalBERT was evaluated on the full "
    "test split (2,926 samples); MedGemma only on a 10% qualitative subset "
    "due to slow autoregressive inference.\n"
    "Direct F1 comparison is not statistically clean.",
    Inches(6.8), Inches(2.38), Inches(6.0), Inches(1.1), ORANGE, ORANGE_LIGHT, 11)

section_label(sl, "Root Cause of Low Overall Scores", Inches(6.8), Inches(3.6))
callout(sl,
    "MedMentions uses a zero-shot label protocol — some test entity types "
    "are entirely absent from training data. This structurally limits the "
    "ceiling for any fine-tuned model regardless of architecture.",
    Inches(6.8), Inches(3.93), Inches(6.0), Inches(0.88), RED, RED_LIGHT, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — QA RESULTS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Results — Task 2: Medical Question Answering")
slide_footer(sl, 9)

metric_box(sl, "25%",  "Random baseline\n(4-choice)",           Inches(0.5),  Inches(0.95), Inches(2.9), Inches(0.9), GRAY)
metric_box(sl, "~26%", "BioClinicalBERT\nTest Accuracy",        Inches(3.55), Inches(0.95), Inches(2.9), Inches(0.9), RED)
metric_box(sl, "55.7%","MedGemma-4B\nTest Accuracy (fine-tuned)",Inches(6.6), Inches(0.95), Inches(2.9), Inches(0.9), ORANGE)
metric_box(sl, "64.4%","MedGemma-4B\nOfficial benchmark",       Inches(9.65), Inches(0.95), Inches(2.9), Inches(0.9), GREEN)

section_label(sl, "Why BioClinicalBERT Underperforms on QA", Inches(0.5), Inches(2.05))
bullet_list(sl, [
    "Max sequence length 128 tokens — USMLE questions often exceed this",
    "Truncation discards part of the question context",
    "Bidirectional encoder not designed for multi-step reasoning",
    "Result: 26% accuracy ≈ 25% random baseline",
], Inches(0.5), Inches(2.38), Inches(5.8), Inches(1.5))

callout(sl,
    "BioClinicalBERT is architecturally ill-suited for QA. "
    "The result confirms that task–architecture mismatch is the dominant factor — "
    "not dataset difficulty.",
    Inches(0.5), Inches(3.98), Inches(5.8), Inches(0.65), RED, RED_LIGHT, 11)

section_label(sl, "MedGemma-4B Fine-tuning Analysis", Inches(6.8), Inches(2.05))
bullet_list(sl, [
    "55.7% test accuracy — 2× random baseline",
    "Official pre-trained benchmark: 64.4% on MedQA",
    "Fine-tuning DEGRADED performance vs. pre-trained baseline",
    "Possible causes: only 2 epochs, LR too high (2e-4), 10 warmup steps",
    "Catastrophic forgetting: fine-tuning may damage medical priors",
], Inches(6.8), Inches(2.38), Inches(6.0), Inches(1.65))

callout(sl,
    "Lesson: a strong pre-trained medical model already 'knows medicine'. "
    "Aggressive fine-tuning on a small dataset can hurt more than it helps. "
    "Conservative instruction tuning (low LR, more warmup) would be safer.",
    Inches(6.8), Inches(4.13), Inches(6.0), Inches(0.78), ORANGE, ORANGE_LIGHT, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — COMPARATIVE ANALYSIS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Comparative Analysis: Encoder vs. Decoder")
slide_footer(sl, 10)

add_table(sl,
    ["Dimension", "Bio_ClinicalBERT (Encoder)", "MedGemma-4B (Decoder)"],
    [["NER Macro F1",     "0.65 (full test, token classifier)", "0.27 (subset, generative)"],
     ["QA Accuracy",      "~26%  (≈ random)",                   "55.7%"],
     ["GPU time / task",  "~20–40 min (110 M params)",           "~1.5–2.5 h (4B, QLoRA)"],
     ["Trainable params", "110 M (full fine-tune)",              "33 M (0.76% via LoRA)"],
     ["Inference speed",  "Fast (single forward pass)",          "Slow (autoregressive)"],
     ["Best suited for",  "Classification, token labeling, NER", "Reasoning, generation, QA"]],
    Inches(0.5), Inches(0.95), W - Inches(1.0), Inches(2.8),
    col_widths=[Inches(2.2), Inches(4.8), Inches(4.8)])

callout(sl,
    "✓  Key takeaway: Use the right architecture for the task.\n"
    "Encoders excel at structured extraction (NER). Decoders excel at open-ended reasoning (QA).\n"
    "Neither model is universally better — task–architecture alignment is the dominant factor.",
    Inches(0.5), Inches(3.92), W - Inches(1.0), Inches(0.88), GREEN, GREEN_LIGHT, 13)

callout(sl,
    "Note: The NER comparison is not a clean head-to-head — BERT used the full "
    "2,926-sample test split while MedGemma used only a 10% qualitative subset. "
    "A fair comparison would evaluate both on the same subset.",
    Inches(0.5), Inches(4.9), W - Inches(1.0), Inches(0.68), ORANGE, ORANGE_LIGHT, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — EFFORT & COMPUTE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Implementation Effort & Compute Estimates")
slide_footer(sl, 11)

section_label(sl, "Development Time (estimated per component)", Inches(0.5), Inches(0.95))
add_table(sl,
    ["Component", "Hours"],
    [["EDA & data pipeline",        "1.5 – 2 h"],
     ["BERT NER baseline",           "4 – 5 h"],
     ["BERT NER with class weights", "5 – 6 h"],
     ["MedGemma NER (QLoRA)",        "6 – 8 h"],
     ["BERT QA",                     "2.5 – 3.5 h"],
     ["MedGemma QA (QLoRA)",         "3 – 4 h"],
     ["Comparative analysis",        "3 – 4 h"],
     ["Report & documentation",      "~10 h"],
     ["Total (estimate)",            "~36 – 44 h"]],
    Inches(0.5), Inches(1.28), Inches(5.8), Inches(3.4),
    col_widths=[Inches(3.8), Inches(2.0)])

section_label(sl, "GPU Execution Time (Google Colab A100)", Inches(6.8), Inches(0.95))
add_table(sl,
    ["Notebook", "Actual / Estimated"],
    [["BERT NER baseline (early stop ~6 ep)",  "~1 – 2 h"],
     ["BERT NER weighted (early stop)",         "~1 – 2 h"],
     ["MedGemma NER (300 steps only)",          "~3 min ⚠"],
     ["BERT QA (early stop ~7 ep)",             "~10 – 15 min"],
     ["MedGemma QA (2 epochs)",                 "~1.5 – 2.5 h"],
     ["Comparative analysis (inference only)",  "~20 – 30 min"],
     ["Total GPU time",                         "~5 – 8 h"]],
    Inches(6.8), Inches(1.28), Inches(6.0), Inches(2.75),
    col_widths=[Inches(3.8), Inches(2.2)])

callout(sl,
    "MedGemma NER trained only 300 of ~1,462 steps per epoch — effectively "
    "ran for 3 minutes. Full training would require 6–12 h. Results are not comparable.",
    Inches(6.8), Inches(4.12), Inches(6.0), Inches(0.65), ORANGE, ORANGE_LIGHT, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — WHAT WORKED / DIDN'T
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Discussion: What Worked & What Didn't")
slide_footer(sl, 12)

section_label(sl, "What Worked", Inches(0.5), Inches(0.95), GREEN)
good = [
    "✓  Class weighting significantly improved NER on frequent entities",
    "✓  QLoRA made the 4B model trainable on a single GPU",
    "✓  MedGemma outperformed BERT 2:1 on QA (55.7% vs 26%)",
    "✓  Custom WeightedTrainer was clean, reusable, and well-documented",
    "✓  Cosine LR + early stopping stabilized training",
    "✓  ECE metric added calibration insight beyond raw accuracy",
    "✓  All models published to HuggingFace Hub for full reproducibility",
    "✓  Unsloth gave 2× training speedup at no additional cost",
]
bullet_list(sl, good, Inches(0.5), Inches(1.28), Inches(6.0), Inches(3.4), size=12)

section_label(sl, "What Didn't Work", Inches(6.8), Inches(0.95), RED)
bad = [
    "✗  CRF layer: fails on 1-line bug in post_init() — never ran",
    "✗  Synthetic data augmentation: 33 samples generated, no impact",
    "✗  BERT QA: near-random accuracy — wrong architecture for reasoning",
    "✗  MedGemma NER: 300 steps only, no quantitative metrics",
    "✗  Fine-tuning MedGemma QA degraded vs pre-trained baseline",
    "✗  Asymmetric NER evaluation — not a clean comparison",
    "✗  No random seeds in most BERT notebooks",
    "✗  Summarization task originally proposed — never implemented",
]
bullet_list(sl, bad, Inches(6.8), Inches(1.28), Inches(6.0), Inches(3.4), size=12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — LIMITATIONS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Limitations")
slide_footer(sl, 13)

section_label(sl, "Data & Evaluation", Inches(0.5), Inches(0.95))
bullet_list(sl, [
    "MedMentions zero-shot label protocol: test entities unseen during training",
    "MedGemma NER evaluated only qualitatively on 3 samples — no reliable F1",
    "BERT QA truncates to 128 tokens — full context lost for longer questions",
    "MedQA is US-centric (USMLE) — does not reflect global clinical practices",
    "MedMentions-MTI881-NER dataset has no stated license — legal status unclear",
], Inches(0.5), Inches(1.28), Inches(5.8), Inches(2.2))

callout(sl,
    "The zero-shot label design is the biggest structural barrier — "
    "even a perfect model cannot predict entity types it has never seen.",
    Inches(0.5), Inches(3.58), Inches(5.8), Inches(0.6), RED, RED_LIGHT, 11)

section_label(sl, "Model & Compute", Inches(6.8), Inches(0.95))
bullet_list(sl, [
    "Single GPU (A100) limits experiment scale and ablation studies",
    "No hyperparameter search — all settings are single-run choices",
    "MedGemma NER severely under-trained (300 of ~1,462 steps per epoch)",
    "No confidence intervals or multi-seed averaging",
    "Catastrophic forgetting risk from fine-tuning MedGemma",
    "BioClinicalBERT trained on MIMIC-III only — US ICU bias",
], Inches(6.8), Inches(1.28), Inches(6.0), Inches(2.2))

callout(sl,
    "Fine-tuning a strong pre-trained medical model requires more care "
    "than fine-tuning a general model — the model already 'knows medicine'.",
    Inches(6.8), Inches(3.58), Inches(6.0), Inches(0.6), ORANGE, ORANGE_LIGHT, 11)

callout(sl,
    "Summary: results are directionally valid but not publication-quality. "
    "The comparative conclusion (encoders for NER, decoders for QA) is robust.",
    Inches(0.5), Inches(4.38), W - Inches(1.0), Inches(0.6), BLUE, BLUE_LIGHT, 12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — FUTURE WORK
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Future Work")
slide_footer(sl, 14)

section_label(sl, "Near-term (straightforward)", Inches(0.5), Inches(0.95))
bullet_list(sl, [
    "Fix CRF notebook (_tied_weights_keys = {}) and run full training",
    "Add random seeds to all BERT notebooks for reproducibility",
    "Run MedGemma NER for 3 full epochs with quantitative seqeval F1",
    "Add constrained decoding for QA (force output to A/B/C/D)",
    "Evaluate BERT QA with Longformer (longer context window)",
    "Add confusion matrix and per-class breakdown to all result tables",
], Inches(0.5), Inches(1.28), Inches(5.8), Inches(2.5))

section_label(sl, "Medium-term (research directions)", Inches(6.8), Inches(0.95))
bullet_list(sl, [
    "RAG: retrieve UMLS definitions to augment NER context at inference",
    "Chain-of-thought prompting for QA reasoning steps",
    "Higher LoRA rank (32–64) + structured output (JSON schema)",
    "Medical summarization task (originally proposed, not implemented)",
    "PubMedBERT as third NER model (originally planned, dropped)",
    "Test on MIMIC-IV notes for real-world deployment evaluation",
    "Build lightweight agent combining NER + QA pipelines",
], Inches(6.8), Inches(1.28), Inches(6.0), Inches(2.5))

callout(sl,
    "The most impactful single improvement: run MedGemma NER for a full epoch "
    "with quantitative evaluation — this would give a valid head-to-head comparison.",
    Inches(0.5), Inches(3.9), W - Inches(1.0), Inches(0.62), BLUE, BLUE_LIGHT, 12)

callout(sl,
    "Consider using the pre-trained MedGemma (no fine-tuning) for QA with "
    "few-shot prompting — may outperform the fine-tuned version.",
    Inches(0.5), Inches(4.6), W - Inches(1.0), Inches(0.62), GREEN, GREEN_LIGHT, 12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — RESPONSIBLE USE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Responsible Use & Ethical Considerations")
slide_footer(sl, 15)

callout(sl,
    "⚠  THESE MODELS MUST NOT BE USED FOR CLINICAL DECISIONS.  "
    "All outputs require review by a qualified medical professional.",
    Inches(0.5), Inches(0.95), W - Inches(1.0), Inches(0.58), RED, RED_LIGHT, 13)

section_label(sl, "Clinical Safety", Inches(0.5), Inches(1.68))
bullet_list(sl, [
    "Models trained on limited datasets — real patients have far more complex presentations",
    "MedGemma can hallucinate plausible but incorrect clinical facts with high confidence",
    "NER errors (false positives/negatives) could misclassify diagnoses or medications",
    "BioClinicalBERT pre-trained on US ICU data only — potential bias toward US populations",
    "55.7% QA accuracy means the model is wrong ~44% of the time — far below clinical threshold",
], Inches(0.5), Inches(2.0), Inches(5.8), Inches(2.3))

section_label(sl, "Data & Licensing", Inches(6.8), Inches(1.68))
bullet_list(sl, [
    "MedMentions-MTI881-NER: no license stated — legally ambiguous",
    "MedQA-USMLE: CC-BY-4.0 — permissive with attribution",
    "MedGemma: Google Health AI Developer Foundations terms — restricts commercial use",
    "MIMIC-III (BioClinicalBERT pre-training): PhysioNet credentialed — not re-distributable",
    "Published HF checkpoints should carry explicit responsible-use disclaimers",
], Inches(6.8), Inches(2.0), Inches(6.0), Inches(2.3))

callout(sl,
    "All code and models are released for educational purposes only. "
    "This project is a technology demonstration, not a clinical tool.",
    Inches(0.5), Inches(4.42), W - Inches(1.0), Inches(0.58), BLUE, BLUE_LIGHT, 12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — THANK YOU / VIDEO
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
add_rect(sl, 0, 0, W, H, fill=RGBColor(0xF0, 0xF6, 0xFF))
add_rect(sl, 0, 0, W, Inches(0.10), fill=BLUE)

add_text_box(sl, "CSCI E-222 · Foundations of Large Language Models · Final Project",
             Inches(0.8), Inches(0.55), W - Inches(1.6), Inches(0.32),
             size=11, color=GRAY)

add_text_box(sl, "Thank You", Inches(0.8), Inches(1.1), W - Inches(1.6), Inches(0.8),
             size=36, bold=True, color=BLUE)

add_text_box(sl, "NLP for Healthcare: Medical Text Analysis",
             Inches(0.8), Inches(1.95), W - Inches(1.6), Inches(0.4),
             size=16, color=GRAY)

# Two info boxes
callout(sl, "Code & Models\nhuggingface.co/alexd063\nCheckpoints for all 4 fine-tuned models",
        Inches(0.8), Inches(2.65), Inches(5.5), Inches(0.88), BLUE, BLUE_LIGHT, 12)
callout(sl, "Datasets Used\nBen10x/MedMentions-MTI881-NER\nGBaker/MedQA-USMLE-4-options",
        Inches(6.6), Inches(2.65), Inches(5.9), Inches(0.88), GREEN, GREEN_LIGHT, 12)

# Video URL box
add_rect(sl, Inches(0.8), Inches(3.78), W - Inches(1.6), Inches(1.35),
         fill=RGBColor(0xFF, 0xF3, 0xCD), line=RGBColor(0xE6, 0xA8, 0x00), line_w=Pt(2))
add_text_box(sl, "Video Presentation URL:", Inches(1.0), Inches(3.9), W - Inches(2.0), Inches(0.3),
             size=12, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
add_text_box(sl, "[ INSERT VIDEO URL HERE ]", Inches(1.0), Inches(4.22), W - Inches(2.0), Inches(0.45),
             size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text_box(sl, "Required: 7–15 minute video covering problem, models, demo, results, and takeaways",
             Inches(1.0), Inches(4.72), W - Inches(2.0), Inches(0.28),
             size=10, color=GRAY, align=PP_ALIGN.CENTER)

add_text_box(sl, "A. Mykhailovskyi  ·  Harvard Extension School  ·  CSCI E-222  ·  May 2026",
             Inches(0.8), H - Inches(0.6), W - Inches(1.6), Inches(0.32),
             size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────────────────
out = "/Users/a.mykhailovskyi/repos/e222/final-project/NLP_Healthcare_Mykhailovskyi.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
